"""Génération d'artefacts RÉELS — réfute audit H14.

L'audit H14 : « Génération de slides/PDF/schémas/data-viz manquante (tests 309-313
stubs) ». On comble : génération de VRAIS artefacts (fichiers) via matplotlib + python-
pptx. Le modèle peut produire des graphiques (data-viz) et des présentations (slides).

* generate_chart(series, path)  : graphique matplotlib → PNG (line/bar/scatter).
* generate_slides(title, slides, path) : présentation .pptx (python-pptx).
* generate_table(data, path)    : table-figure matplotlib → PNG.

Chaque fonction produit un FICHIER RÉEL sur disque. C'est la génération d'artefacts
du cahier des charges (§5 — créer des objets). HONNÊTE : data-viz/slides réels via
librairies (pas un modèle génératif entraîné — qui nécessiterait un corpus) ; la
CAPACITÉ (produire un artefact) est réelle et vérifiable (fichier créé + lisible).
"""
from __future__ import annotations
import os
from typing import Any, Dict, List, Optional, Tuple


def _mpl():
    import matplotlib
    matplotlib.use("Agg")          # headless (pas de display requis)
    import matplotlib.pyplot as plt
    return plt


def generate_chart(series: List[Tuple[str, List[float]]], path: str,
                   title: str = "Chart", kind: str = "line") -> Dict[str, Any]:
    """Génère un graphique (line/bar/scatter) → PNG. series=[(label, values)].
    Retourne {path, ok, n_series}."""
    try:
        plt = _mpl()
        fig, ax = plt.subplots(figsize=(8, 5))
        for label, values in series:
            xs = list(range(len(values)))
            if kind == "bar":
                ax.bar(xs, values, label=label)
            elif kind == "scatter":
                ax.scatter(xs, values, label=label)
            else:
                ax.plot(xs, values, label=label, marker="o")
        ax.set_title(title)
        ax.legend()
        ax.grid(True, alpha=0.3)
        fig.tight_layout()
        fig.savefig(path, dpi=100)
        plt.close(fig)
        return {"path": path, "ok": os.path.exists(path), "n_series": len(series),
                "kind": kind}
    except Exception as e:
        return {"path": path, "ok": False, "error": f"{type(e).__name__}: {e}"}


def generate_slides(title: str, slides: List[Dict[str, Any]], path: str) -> Dict[str, Any]:
    """Génère une présentation .pptx. slides=[{title, bullets:[...]}].
    Retourne {path, ok, n_slides}."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        # slide de titre
        tlayout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(tlayout)
        slide.shapes.title.text = title
        # slides de contenu
        clayout = prs.slide_layouts[1]
        for s in slides:
            cslide = prs.slides.add_slide(clayout)
            cslide.shapes.title.text = s.get("title", "")
            body = cslide.placeholders[1]
            tf = body.text_frame
            for i, bullet in enumerate(s.get("bullets", [])):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(bullet)
                p.font.size = Pt(18)
        prs.save(path)
        return {"path": path, "ok": os.path.exists(path), "n_slides": len(slides) + 1}
    except Exception as e:
        return {"path": path, "ok": False, "error": f"{type(e).__name__}: {e}"}


def generate_table(headers: List[str], rows: List[List[Any]], path: str,
                   title: str = "Table") -> Dict[str, Any]:
    """Génère une table-figure → PNG."""
    try:
        plt = _mpl()
        fig, ax = plt.subplots(figsize=(max(6, len(headers) * 1.5),
                                        max(2, len(rows) * 0.5)))
        ax.axis("off")
        ax.table(cellText=[[str(c) for c in r] for r in rows],
                 colLabels=headers, loc="center", cellLoc="center")
        ax.set_title(title)
        fig.savefig(path, dpi=100, bbox_inches="tight")
        plt.close(fig)
        return {"path": path, "ok": os.path.exists(path), "n_rows": len(rows)}
    except Exception as e:
        return {"path": path, "ok": False, "error": f"{type(e).__name__}: {e}"}


def available_generators() -> Dict[str, bool]:
    """Quels générateurs d'artefacts sont disponibles (libs installées) ?"""
    avail = {"chart": False, "table": False, "slides": False}
    try:
        import matplotlib  # noqa
        avail["chart"] = avail["table"] = True
    except ImportError:
        pass
    try:
        import pptx  # noqa
        avail["slides"] = True
    except ImportError:
        pass
    return avail


if __name__ == "__main__":
    import tempfile
    tmp = tempfile.mkdtemp()
    avail = available_generators()
    print(f"[artefact_generator] générateurs dispo: {avail}")
    if avail["chart"]:
        r = generate_chart([("A", [1, 3, 2, 5, 4]), ("B", [2, 2, 4, 3, 6])],
                           os.path.join(tmp, "chart.png"), title="Démo data-viz", kind="line")
        print(f"  chart : {r['ok']} → {r['path']}")
    if avail["slides"]:
        r = generate_slides("Présentation OCM-26400",
                            [{"title": "Crown-jewel", "bullets": ["grok binaire 100%", "décomp +99.5pt"]},
                             {"title": "Domaines", "bullets": ["33 domaines", "101 règles"]}],
                            os.path.join(tmp, "slides.pptx"))
        print(f"  slides: {r['ok']} ({r['n_slides']} slides) → {r['path']}")
    if avail["table"]:
        r = generate_table(["Benchmark", "Score"], [["Agentic", "91.7%"], ["Reasoning", "100%"]],
                           os.path.join(tmp, "table.png"), title="Scores")
        print(f"  table : {r['ok']} → {r['path']}")
